from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_final_accurate_presentation():
    prs = Presentation()
    
    # Image Paths
    title_bg = r"C:\Users\HP\.gemini\antigravity\brain\642ee762-76c8-4b4a-8247-c3e87c808fbd\tendershield_title_bg_1778074469059.png"
    internal_bg = r"C:\Users\HP\.gemini\antigravity\brain\642ee762-76c8-4b4a-8247-c3e87c808fbd\tendershield_dark_internal_bg_v2_1778074752755.png"

    # Colors
    BLUE = RGBColor(56, 189, 248)
    WHITE = RGBColor(255, 255, 255)
    NAVY = RGBColor(10, 25, 47)

    def set_bg(slide, image_path):
        slide.shapes.add_picture(image_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    def add_premium_header(slide, title_text):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = NAVY
        shape.line.color.rgb = BLUE
        shape.line.width = Pt(1.5)
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(8.6), Inches(0.7))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = WHITE

    def add_slide(title_text, points):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide, internal_bg)
        add_premium_header(slide, title_text)
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5))
        tf = content_box.text_frame
        tf.word_wrap = True
        for i, point in enumerate(points):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = f"▶  {point}"
            p.font.size = Pt(18)
            p.font.color.rgb = WHITE
            p.space_after = Pt(14)

    # 1. TITLE SLIDE (FIXED: Standard Layout to prevent overlap)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_bg(slide, title_bg)
    # Bring title and subtitle to front
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "TenderShield AI"
    title.text_frame.paragraphs[0].font.size = Pt(60)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = WHITE
    
    subtitle.text = "Automated Evaluation & Fraud Detection System\nCRPF PROCUREMENT MISSION CONTROL"
    subtitle.text_frame.paragraphs[0].font.size = Pt(20)
    subtitle.text_frame.paragraphs[0].font.color.rgb = BLUE

    # 2. Problem
    add_slide("The Challenge", [
        "Inefficient Manual Audits: Reviewing 100+ pages per tender is prone to human error.",
        "The 'Black Box' Barrier: Traditional AI lacks the transparency required for government use.",
        "Systemic Fraud: Cartels and shell companies bypass static manual filters.",
        "Compliance Risk: Missing a single mandatory document can lead to legal litigation."
    ])

    # 3. Solution
    add_slide("TenderShield: Mission Control", [
        "End-to-End Automation: From PDF upload to L1 Ranking.",
        "Neural Rule Extraction: LLM-powered parsing of mandatory legal & financial rules.",
        "Explainable AI: Every decision is backed by a verifiable 'Drill-Down' matrix.",
        "Persistence & Integrity: Locally persisted database with safe archival storage."
    ])

    # 4. Secure Gateway
    add_slide("Secure Officer Gateway", [
        "Credential Protected: Access restricted to authorized CRPF procurement officers.",
        "Session Persistence: Maintaining state across browser refreshes for seamless demos.",
        "Enterprise Dashboard: High-fidelity, reactive UI built for speed and security.",
        "System Transparency: Glowing status badges indicating neural engine health."
    ])

    # 5. Document AI
    add_slide("Neural Rule Extraction", [
        "Intelligent PDF Parsing: Automatically extracting thresholds from Notice Inviting Tenders (NIT).",
        "Semantic Clause Analysis: AI identifies 'Mandatory' vs. 'Optional' legal rules.",
        "Verifiable Sources: AI captures the exact clause text for officer verification.",
        "Category Logic: Rules are automatically mapped to Financial, Technical, and Legal buckets."
    ])

    # 6. AI Bidder Entry
    add_slide("AI-Assisted Bidder Entry", [
        "LLM Auto-Fill: Upload bidder documents to automatically populate company profiles.",
        "Data Extraction: Automatically identifies Turnover, Experience, and Legal IDs.",
        "Manual Override: Officers can refine AI-extracted bidder values in real-time.",
        "Compliance Checkboxes: Integrated legal document verification (PAN, GST, License)."
    ])

    # 7. Drill-Down Matrix
    add_slide("Explainable Evaluation Matrix", [
        "Zero Silent Rejections: Full transparency into every Pass/Fail decision.",
        "The Matrix View: Granular breakdown of Required vs. Actual bidder parameters.",
        "Drill-Down Logic: Click any bidder to see exactly which rules they failed.",
        "Officer Control: AI assists, but the officer always has the final oversight."
    ])

    # 8. Fraud Intelligence
    add_slide("Threat Intelligence Dashboard", [
        "Anomaly Detection: Real-time identification of predatory bidding and cartels.",
        "Dynamic Confidence: AI-calculated probability scores for every security flag.",
        "Severity Tiering: Instant flagging of 🛑 High Risk vs ⚠️ Suspicious bidders.",
        "Evidence Log: Detailed descriptions of detected inconsistencies in bidder data."
    ])

    # 9. Vigilance Engine
    add_slide("Vigilance: Cartel & Collusion", [
        "Network Analysis: Automatic detection of shared Directors across competing bids.",
        "IP Identification: Flagging multiple submissions originating from the same subnet.",
        "Compliance Mapping: Identifying legal non-compliance and missing credentials.",
        "Blacklist Linkage: Cross-referencing bidders against debarment databases."
    ])

    # 10. Leaderboard Mechanics
    add_slide("Strategic L1 Ranking", [
        "Dynamic Scoring: Composite weighted scores based on technical and financial strength.",
        "Bonus Multipliers: Scoring bidders who exceed mandatory tender requirements.",
        "Real-Time Charts: Visual bar charts highlighting the 'Strategic L1 Winner' in Gold.",
        "Risk Penalization: Automatically adjusting rankings based on fraud risk levels."
    ])

    # 11. Project Archives
    add_slide("Persistent Audit & Archives", [
        "Audit Transparency: Historical logging of every evaluation performed.",
        "Project Archives: Safe, persistent storage of past tender outcomes.",
        "Cascading Deletion: Secure right-click 'Delete' to scrub records from the database.",
        "Stateful Recovery: Database-driven state ensures zero data loss during demos."
    ])

    # 12. Final Impact
    add_slide("Mission Ready", [
        "TenderShield AI: Protecting the Integrity of Procurement.",
        "Efficiency: 10x faster screening of mandatory bidder eligibility.",
        "Trust: 100% auditability through Explainable AI matrices.",
        "Ready for Submission: A fully functional prototype for CRPF evaluation."
    ])

    prs.save('TenderShield_AI_Final_Submission_Deck.pptx')
    print("Final Accurate 12-Slide Deck saved as TenderShield_AI_Final_Submission_Deck.pptx")

if __name__ == "__main__":
    create_final_accurate_presentation()
